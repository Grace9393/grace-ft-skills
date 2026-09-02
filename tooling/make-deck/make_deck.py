#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""make_deck.py - one command from a source file to a finished IBM-branded deck.

Runs the file-to-pptx converter with the bundled IBM Carbon template, then
restores the Human / AI / Hybrid colour semantics that live in the source HTML's
CSS classes (pill-h / pill-a / pill-ha / pill-n) and are lost by any text
extraction. Both steps are the same code the skill runs; this just chains them
so the result needs no manual pass.

Usage
    python make_deck.py INPUT [INPUT ...] -o OUTPUT.pptx
    python make_deck.py report.html                 # -> report.pptx beside it

Options
    -o, --output   output .pptx (default: first input's name with .pptx)
    --no-shading   skip the colour-semantics restoration
    --skill DIR    skill bundle to use (default: auto-detect)

Requires: python-pptx, openpyxl, pandas, pdfplumber, pillow, lxml
    pip install python-pptx openpyxl pandas pdfplumber pillow lxml
"""
from __future__ import annotations

import argparse
import json
import re
import subprocess
import sys
from collections import defaultdict
from pathlib import Path

# Candidate locations for the file-to-pptx bundle, best first.
SKILL_CANDIDATES = [
    Path(__file__).resolve().parents[2] / "skills" / "file-to-pptx",
    Path(r"C:\Users\GRACEPAN\.bob\skills\file-to-pptx"),
    Path(r"H:\My Drive\AA\.bob\skills\file-to-pptx"),
    Path(r"C:\Users\GRACEPAN\.claude\skills\file-to-pptx"),
]

SCHEME = {
    "pill-h":  ((0xE0, 0xE0, 0xE0), "Human"),
    "pill-a":  ((0xD0, 0xE2, 0xFF), "AI"),
    "pill-ha": ((0xE8, 0xDA, 0xFF), "Hybrid"),
    "pill-n":  ((0xF4, 0xF4, 0xF4), "N/A"),
}
ORDER = ["pill-h", "pill-ha", "pill-a", "pill-n"]


def find_skill(explicit: str | None) -> Path:
    if explicit:
        p = Path(explicit)
        if (p / "scripts" / "universal_to_pptx.py").is_file():
            return p
        sys.exit(f"no converter under {p}")
    for p in SKILL_CANDIDATES:
        if (p / "scripts" / "universal_to_pptx.py").is_file():
            return p
    sys.exit("file-to-pptx bundle not found; pass --skill DIR")


def convert(skill: Path, inputs: list[str], out: Path) -> dict:
    cmd = [sys.executable, str(skill / "scripts" / "universal_to_pptx.py"),
           *inputs, "-o", str(out)]
    tpl, brand = skill / "assets" / "ibm_template.pptx", skill / "assets" / "ibm_carbon.json"
    if tpl.is_file():
        cmd += ["--template", str(tpl)]
    if brand.is_file():
        cmd += ["--brand", str(brand)]
    r = subprocess.run(cmd, capture_output=True, text=True, encoding="utf-8",
                       errors="replace")
    if r.returncode != 0:
        sys.stderr.write(r.stdout + r.stderr)
        sys.exit(f"converter failed ({r.returncode})")
    m = re.search(r"\{.*\}", r.stdout, re.S)
    return json.loads(m.group(0)) if m else {}


def shade(pptx: Path, html_sources: list[Path]) -> tuple[int, dict, list[int]]:
    """Colour table cells whose text came from a classed pill in the source."""
    try:
        from lxml import html as lhtml
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return 0, {}, []

    mapping, seen = {}, defaultdict(set)
    for src in html_sources:
        try:
            tree = lhtml.fromstring(src.read_text(encoding="utf-8", errors="ignore"))
        except Exception:
            continue
        for el in tree.iter():
            cls = (el.get("class") or "").split()
            if "pill" not in cls:
                continue
            kind = next((c for c in cls if c in SCHEME), None)
            if not kind:
                continue
            txt = " ".join("".join(el.itertext()).split())
            if txt:
                seen[txt].add(kind)
    for txt, kinds in seen.items():
        if len(kinds) == 1:                       # skip anything ambiguous
            mapping[txt] = next(iter(kinds))
    if not mapping:
        return 0, {}, []

    prs = Presentation(str(pptx))
    n, per, slides = 0, defaultdict(int), set()
    for idx, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            if not shp.has_table:
                continue
            for row in shp.table.rows:
                for cell in row.cells:
                    k = mapping.get(" ".join(cell.text.split()))
                    if not k:
                        continue
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*SCHEME[k][0])
                    n += 1
                    per[k] += 1
                    slides.add(idx)
    for idx in sorted(slides):                    # legend
        s = prs.slides[idx - 1]
        box = s.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.85),
                                   prs.slide_width - Inches(1), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        box.text_frame.word_wrap = False
        r0 = p.add_run(); r0.text = "Responsibility:  "
        r0.font.size = Pt(11); r0.font.bold = True
        for k in ORDER:
            if per.get(k):
                r = p.add_run(); r.text = f"  \u25a0 {SCHEME[k][1]}  "
                r.font.size = Pt(11)
                r.font.color.rgb = RGBColor(*(max(0, c - 0x50) for c in SCHEME[k][0]))
    prs.save(str(pptx))
    return n, dict(per), sorted(slides)


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("inputs", nargs="+")
    ap.add_argument("-o", "--output")
    ap.add_argument("--no-shading", action="store_true")
    ap.add_argument("--skill")
    a = ap.parse_args()

    for f in a.inputs:
        if not Path(f).is_file():
            sys.exit(f"no such file: {f}")
    out = Path(a.output) if a.output else Path(a.inputs[0]).with_suffix(".pptx")
    skill = find_skill(a.skill)

    print(f"bundle : {skill}")
    print(f"input  : {', '.join(a.inputs)}")
    summary = convert(skill, a.inputs, out)
    print(f"slides : {summary.get('slide_count', '?')}")
    for w in summary.get("warnings", []):
        print(f"  warning: {w}")

    if not a.no_shading:
        htmls = [Path(f) for f in a.inputs if Path(f).suffix.lower() in (".html", ".htm")]
        if htmls:
            n, per, slides = shade(out, htmls)
            if n:
                print("shading: " + ", ".join(f"{SCHEME[k][1]} {v}"
                                              for k, v in per.items())
                      + f"  (slides {slides})")
            else:
                print("shading: no classed pills in source - nothing to restore")

    print(f"\nwrote  : {out}  ({out.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
